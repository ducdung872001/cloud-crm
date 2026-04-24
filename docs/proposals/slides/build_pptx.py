"""
Build TNTech branded PPTX from markdown slide decks.
Brand colors derived from tntech logo: cyan -> blue gradient.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw, ImageFont

# ===== BRAND COLORS (từ logo TNTech) =====
BRAND_CYAN = RGBColor(0x1E, 0xC1, 0xF5)      # #1EC1F5 - cyan sáng
BRAND_BLUE = RGBColor(0x1E, 0x9F, 0xE5)      # #1E9FE5 - xanh TNTech chính
BRAND_DARK = RGBColor(0x0E, 0x4F, 0x8C)      # #0E4F8C - xanh đậm chữ 'e'
BRAND_DARKER = RGBColor(0x07, 0x2D, 0x56)    # #072D56 - nền slide tối
BRAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_OFFWHITE = RGBColor(0xF4, 0xF9, 0xFD)  # nền content nhạt
BRAND_GRAY = RGBColor(0x4A, 0x5A, 0x6A)
BRAND_LIGHT = RGBColor(0xE5, 0xF2, 0xFB)     # nhẹ nhẹ cho alternate row
BRAND_ORANGE = RGBColor(0xF5, 0x8F, 0x1E)    # accent warm
BRAND_YELLOW = RGBColor(0xFF, 0xC1, 0x07)    # accent highlight

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSETS_DIR = Path(__file__).parent / "assets"
ASSETS_DIR.mkdir(exist_ok=True)
LOGO_PATH = ASSETS_DIR / "tntech-logo.png"
LOGO_WHITE_PATH = ASSETS_DIR / "tntech-logo-white.png"


# ============================================================
# 1. LOGO GENERATION (mô phỏng logo tntech gradient cyan→blue)
# ============================================================
def make_logo(path: Path, for_dark_bg=False):
    """Render logo 'tntech' mô phỏng logo TNTech:
    - Các chữ 't n t c h' màu cyan sáng
    - Chữ 'e' đặc biệt: vẽ dạng CIRCLE (giống iris/camera aperture) màu xanh đậm,
      có dot nhỏ chính giữa — điểm nhấn signature của TNTech.
    """
    W, H = 1400, 380
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    # Load bold rounded font
    font = None
    for candidate in [
        r"C:\Windows\Fonts\segoeuib.ttf",
        r"C:\Windows\Fonts\arialbd.ttf",
        r"C:\Windows\Fonts\calibrib.ttf",
    ]:
        try:
            font = ImageFont.truetype(candidate, 260)
            break
        except OSError:
            continue
    if font is None:
        font = ImageFont.load_default()

    # Colors
    CYAN = (0x1E, 0xC1, 0xF5, 255)       # chữ sáng
    CYAN_MID = (0x1E, 0xA8, 0xE9, 255)   # biến thể nhẹ
    DARK = (0x0E, 0x4F, 0x8C, 255)       # chữ 'e' + dot

    # Measure each char width
    def cw(ch):
        b = draw.textbbox((0, 0), ch, font=font)
        return b[2] - b[0]

    # Layout: tnt + [CIRCLE 'e'] + ch
    w_t = cw("t")
    w_n = cw("n")
    w_c = cw("c")
    w_h = cw("h")
    spacing = 12
    # Circle 'e' sẽ lớn bằng chiều cao chữ thường × 0.72 × em
    _, desc = font.getmetrics()
    asc, _ = font.getmetrics()
    x_height = int(font.size * 0.54)  # approx x-height for Segoe UI Bold
    circle_d = int(x_height * 1.15)  # hơi lớn hơn x-height cho nổi bật

    total_w = w_t + w_n + w_t + circle_d + w_c + w_h + spacing * 5
    x = (W - total_w) // 2
    # Baseline y so that chars are vertically centered
    y_text = (H - asc) // 2 - 10

    # Draw "t n t"
    draw.text((x, y_text), "t", font=font, fill=CYAN)
    x += w_t + spacing
    draw.text((x, y_text), "n", font=font, fill=CYAN)
    x += w_n + spacing
    draw.text((x, y_text), "t", font=font, fill=CYAN_MID)
    x += w_t + spacing

    # Draw circle 'e' — vertically centered with x-height of other chars
    # Baseline is at y_text + asc, x-height top is at y_text + asc - x_height
    baseline_y = y_text + asc
    xh_top = baseline_y - x_height
    xh_center_y = (xh_top + baseline_y) // 2

    r_outer = circle_d // 2
    cx_e = x + r_outer
    cy_e = xh_center_y

    # Outer filled circle (dark blue)
    draw.ellipse(
        (cx_e - r_outer, cy_e - r_outer, cx_e + r_outer, cy_e + r_outer),
        fill=DARK,
    )
    # Cut-out to form 'e' shape: a notch on the right (where 'e' opens)
    notch_w = int(r_outer * 0.65)
    notch_h = int(r_outer * 0.22)
    draw.rectangle(
        (cx_e, cy_e - notch_h, cx_e + notch_w + 4, cy_e + notch_h),
        fill=(0, 0, 0, 0),
    )
    # Inner solid dot (signature feature)
    r_dot = int(r_outer * 0.26)
    draw.ellipse(
        (cx_e - r_dot, cy_e - r_dot, cx_e + r_dot, cy_e + r_dot),
        fill=DARK,
    )

    x += circle_d + spacing
    # Draw "c h"
    draw.text((x, y_text), "c", font=font, fill=CYAN_MID)
    x += w_c + spacing
    draw.text((x, y_text), "h", font=font, fill=CYAN)

    # Crop tight
    bbox = img.getbbox()
    if bbox:
        img = img.crop((
            max(0, bbox[0] - 20),
            max(0, bbox[1] - 20),
            min(W, bbox[2] + 20),
            min(H, bbox[3] + 20),
        ))
    img.save(path)
    return path


make_logo(LOGO_PATH, for_dark_bg=False)
make_logo(LOGO_WHITE_PATH, for_dark_bg=True)


# ============================================================
# 2. MARKDOWN PARSER
# ============================================================
class Slide:
    def __init__(self):
        self.title = ""
        self.blocks = []  # list of ("bullet"/"para"/"table"/"quote"/"sub", content)


def parse_markdown(md_text: str) -> tuple[dict, list[Slide]]:
    """Parse markdown into metadata + slides split by '## '."""
    meta = {}
    lines = md_text.split("\n")
    i = 0

    # YAML frontmatter
    if lines and lines[0].strip() == "---":
        i = 1
        while i < len(lines) and lines[i].strip() != "---":
            m = re.match(r'^(\w+):\s*"?(.+?)"?\s*$', lines[i])
            if m:
                meta[m.group(1)] = m.group(2)
            i += 1
        i += 1

    slides = []
    current = None
    buf = []

    def flush_buf():
        nonlocal buf
        if not buf:
            return
        # Detect kind
        text = "\n".join(buf).strip()
        if not text:
            buf = []
            return
        # Table?
        if any(l.strip().startswith("|") for l in buf):
            current.blocks.append(("table", buf[:]))
        elif all(l.strip().startswith(("- ", "* ", "1.", "2.", "3.", "4.", "5.",
                                         "6.", "7.", "8.", "9.", "10.")) or not l.strip()
                 for l in buf):
            bullets = [re.sub(r"^(\d+\.|\-|\*)\s*", "", l).strip()
                       for l in buf if l.strip()]
            current.blocks.append(("bullet", bullets))
        elif text.startswith(">"):
            current.blocks.append(("quote",
                                   re.sub(r"^>\s*", "", text, flags=re.M)))
        else:
            current.blocks.append(("para", text))
        buf = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped == "---":
            # slide separator
            if current:
                flush_buf()
                slides.append(current)
                current = None
            i += 1
            continue

        # H1 - cover
        if stripped.startswith("# ") and current is None:
            current = Slide()
            current.title = stripped[2:].strip()
            current._is_cover = True
            i += 1
            continue

        # H2 - new slide
        if stripped.startswith("## "):
            if current:
                flush_buf()
                slides.append(current)
            current = Slide()
            current.title = stripped[3:].strip()
            i += 1
            continue

        if current is None:
            i += 1
            continue

        # Bullet block: consecutive lines starting with '-' or number
        if re.match(r"^\s*([\-\*]|\d+\.)\s", line):
            # Flush non-bullet buffer
            if buf and not all(re.match(r"^\s*([\-\*]|\d+\.)\s", l) or not l.strip()
                               for l in buf):
                flush_buf()
            buf.append(line)
            i += 1
            continue

        # Table block
        if stripped.startswith("|"):
            if buf and not any(l.strip().startswith("|") for l in buf):
                flush_buf()
            buf.append(line)
            i += 1
            continue

        # Quote
        if stripped.startswith(">"):
            if buf:
                flush_buf()
            qbuf = [line]
            i += 1
            while i < len(lines) and lines[i].strip().startswith(">"):
                qbuf.append(lines[i])
                i += 1
            current.blocks.append(
                ("quote", "\n".join(re.sub(r"^>\s*", "", l) for l in qbuf))
            )
            continue

        # Blank line - flush bullets/tables
        if not stripped:
            if buf:
                flush_buf()
            i += 1
            continue

        # Regular paragraph - collect until blank
        if buf and any(l.strip().startswith("|") or
                       re.match(r"^\s*([\-\*]|\d+\.)\s", l) for l in buf):
            flush_buf()
        buf.append(line)
        i += 1

    if current:
        flush_buf()
        slides.append(current)

    return meta, slides


def parse_table(lines: list[str]) -> list[list[str]]:
    """Parse markdown table rows."""
    rows = []
    for l in lines:
        l = l.strip()
        if not l.startswith("|"):
            continue
        # Skip separator row |---|---|
        if re.match(r"^\|?\s*[\-:]+\s*(\|\s*[\-:]+\s*)+\|?$", l):
            continue
        cells = [c.strip() for c in l.strip("|").split("|")]
        rows.append(cells)
    return rows


# ============================================================
# 3. PPTX HELPERS
# ============================================================
def set_slide_bg_color(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_rect(slide, left, top, width, height, colors):
    """Draw a decorative rectangle (solid - simple)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors[0]
    return shape


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="Calibri",
                 anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.name = font_name
    if color is not None:
        r.font.color.rgb = color
    return tb


def add_rich_para(tf, text, font_size=16, color=None, bold=False,
                  bullet=False, font_name="Calibri", level=0):
    """Add a paragraph with inline **bold** parsing."""
    p = tf.add_paragraph() if tf.paragraphs[0].text or tf.paragraphs[0].runs \
        else tf.paragraphs[0]
    p.level = level
    # Remove any existing runs
    p.text = ""

    # Parse **bold** and *italic* inline
    parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)
    for part in parts:
        if not part:
            continue
        is_bold = bold
        is_italic = False
        t = part
        if part.startswith("**") and part.endswith("**"):
            t = part[2:-2]
            is_bold = True
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            t = part[1:-1]
            is_italic = True
        r = p.add_run()
        r.text = t
        r.font.size = Pt(font_size)
        r.font.bold = is_bold
        r.font.italic = is_italic
        r.font.name = font_name
        if color is not None:
            r.font.color.rgb = color

    if bullet:
        # Add bullet via XML
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "●")
    return p


def add_header_bar(slide, title_text, slide_num=None, total_slides=None,
                   is_section=False):
    """Add a decorative header bar with title + accent."""
    # Left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(0.25), SLIDE_H)
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = BRAND_BLUE

    # Top bar
    topbar_h = Inches(0.65)
    topbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, topbar_h)
    topbar.line.fill.background()
    topbar.fill.solid()
    topbar.fill.fore_color.rgb = BRAND_DARKER

    # Accent cyan line under top bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, topbar_h,
                                     SLIDE_W, Inches(0.05))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_CYAN

    # Logo in top-left of bar
    try:
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.35), Inches(0.12),
            height=Inches(0.42)
        )
    except Exception:
        pass

    # Title text (top bar right-aligned text)
    tb = slide.shapes.add_textbox(Inches(2.2), Inches(0.08),
                                    SLIDE_W - Inches(3.5), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = BRAND_WHITE
    r.font.name = "Calibri"

    # Slide number (bottom-right)
    if slide_num is not None:
        nb = slide.shapes.add_textbox(
            SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.45),
            Inches(1.2), Inches(0.3)
        )
        nt = nb.text_frame
        np_ = nt.paragraphs[0]
        np_.alignment = PP_ALIGN.RIGHT
        nr = np_.add_run()
        nr.text = f"{slide_num} / {total_slides}"
        nr.font.size = Pt(10)
        nr.font.color.rgb = BRAND_GRAY

    # Footer brand
    fb = slide.shapes.add_textbox(Inches(0.5), SLIDE_H - Inches(0.45),
                                    Inches(6), Inches(0.3))
    ft = fb.text_frame
    fp = ft.paragraphs[0]
    fr = fp.add_run()
    fr.text = "TNTech — Chuyển đổi số, làm chủ tương lai  |  tntech.vn"
    fr.font.size = Pt(10)
    fr.font.italic = True
    fr.font.color.rgb = BRAND_GRAY


# ============================================================
# 4. SLIDE BUILDERS
# ============================================================
def build_cover_slide(prs, title, subtitle, date_text):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg_color(slide, BRAND_DARKER)

    # Decorative gradient-like bands (simulate gradient with multiple rects)
    for i, (h, color, alpha) in enumerate([
        (Inches(1.3), BRAND_DARK, None),
        (Inches(0.08), BRAND_BLUE, None),
        (Inches(0.04), BRAND_CYAN, None),
    ]):
        y = SLIDE_H - sum([Inches(1.3), Inches(0.08), Inches(0.04)][:i + 1])
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, h)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = color

    # Corner accents
    corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, Inches(2.5), Inches(2.5)
    )
    corner.line.fill.background()
    corner.fill.solid()
    corner.fill.fore_color.rgb = BRAND_BLUE

    # Logo center top
    try:
        slide.shapes.add_picture(
            str(LOGO_WHITE_PATH),
            Inches(5.17), Inches(1.0),
            width=Inches(3.0)
        )
    except Exception:
        pass

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.0),
                                    SLIDE_W - Inches(1.6), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = BRAND_WHITE
    r.font.name = "Calibri"

    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.3),
                                    SLIDE_W - Inches(1.6), Inches(0.9))
    sf = sb.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.size = Pt(24)
    sr.font.color.rgb = BRAND_CYAN
    sr.font.name = "Calibri Light"

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.67), Inches(5.3),
        Inches(2), Inches(0.04)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_CYAN

    # Date / tagline
    db = slide.shapes.add_textbox(Inches(0.8), Inches(5.5),
                                    SLIDE_W - Inches(1.6), Inches(0.5))
    dt = db.text_frame
    dp = dt.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = date_text
    dr.font.size = Pt(16)
    dr.font.italic = True
    dr.font.color.rgb = BRAND_WHITE


def build_thankyou_slide(prs, title, bullets):
    """Closing slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg_color(slide, BRAND_DARKER)

    # Top decorative
    for y, h, c in [
        (0, Inches(0.08), BRAND_CYAN),
        (Inches(0.08), Inches(0.04), BRAND_BLUE),
    ]:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, h)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = c

    # Logo
    try:
        slide.shapes.add_picture(
            str(LOGO_WHITE_PATH),
            Inches(5.17), Inches(0.6),
            width=Inches(3.0)
        )
    except Exception:
        pass

    # Big "Cảm ơn!" centered
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.3),
                                    SLIDE_W - Inches(1.6), Inches(1.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = BRAND_WHITE
    r.font.name = "Calibri"

    # Contact info
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(4.0),
                                    SLIDE_W - Inches(1.6), Inches(2.5))
    cf = cb.text_frame
    cf.word_wrap = True
    for idx, line in enumerate(bullets):
        p = cf.add_paragraph() if idx > 0 else cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18)
        r.font.color.rgb = BRAND_CYAN if "@" in line or "tntech" in line.lower() \
            else BRAND_WHITE
        r.font.name = "Calibri"
        r.font.bold = ("@" in line or "tntech" in line.lower())


def build_agenda_slide(prs, items, slide_num, total_slides):
    """Table-of-contents style slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg_color(slide, BRAND_OFFWHITE)
    add_header_bar(slide, "NỘI DUNG TRÌNH BÀY", slide_num, total_slides)

    # Slide title
    tb = add_text_box(slide, Inches(0.5), Inches(1.0),
                       SLIDE_W - Inches(1.0), Inches(0.6),
                       "Nội dung trình bày", font_size=32, bold=True,
                       color=BRAND_DARK)

    # Two-column layout
    n = len(items)
    mid = (n + 1) // 2
    col_w = Inches(5.8)
    col_y = Inches(1.85)
    col_h = Inches(5.0)

    for col_idx, col_items in enumerate([items[:mid], items[mid:]]):
        x = Inches(0.5) + col_idx * Inches(6.3)
        tb = slide.shapes.add_textbox(x, col_y, col_w, col_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(col_items):
            idx = col_idx * mid + i + 1
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = ""
            p.space_after = Pt(10)
            # number badge
            r1 = p.add_run()
            r1.text = f"{idx:02d}  "
            r1.font.size = Pt(20)
            r1.font.bold = True
            r1.font.color.rgb = BRAND_BLUE
            r1.font.name = "Calibri"
            # item text
            # Strip markdown **x**
            clean = re.sub(r"\*\*([^*]+)\*\*", r"\1", item)
            r2 = p.add_run()
            r2.text = clean
            r2.font.size = Pt(15)
            r2.font.color.rgb = BRAND_DARKER
            r2.font.name = "Calibri"


def build_section_divider(prs, section_num, section_title, slide_num,
                           total_slides):
    """Full-colored section divider slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg_color(slide, BRAND_DARK)

    # Corner triangle
    corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, Inches(2), Inches(2)
    )
    corner.line.fill.background()
    corner.fill.solid()
    corner.fill.fore_color.rgb = BRAND_BLUE

    # Bottom decoration
    for i, (h, c) in enumerate([
        (Inches(0.04), BRAND_CYAN),
        (Inches(0.08), BRAND_BLUE),
    ]):
        y = SLIDE_H - sum([Inches(0.04), Inches(0.08)][:i + 1])
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, h)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = c

    # Big number
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
                                    Inches(4), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{section_num:02d}"
    r.font.size = Pt(160)
    r.font.bold = True
    r.font.color.rgb = BRAND_CYAN
    r.font.name = "Calibri"

    # Section title
    tb = slide.shapes.add_textbox(Inches(5.0), Inches(3.0),
                                    Inches(7.8), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section_title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = BRAND_WHITE
    r.font.name = "Calibri"

    # Small label
    tb = slide.shapes.add_textbox(Inches(5.0), Inches(2.5),
                                    Inches(7), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "PHẦN"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = BRAND_CYAN
    r.font.name = "Calibri"


def clean_md(text: str) -> str:
    """Strip markdown formatting for plain display."""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


def build_content_slide(prs, slide_data: Slide, slide_num, total_slides):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg_color(slide, BRAND_OFFWHITE)

    # Strip leading number "1. " from title for header bar
    header_title = re.sub(r"^\d+\.\s*", "", slide_data.title).upper()
    add_header_bar(slide, header_title, slide_num, total_slides)

    # Slide title (inside main content)
    title_clean = clean_md(slide_data.title)
    tb = add_text_box(
        slide, Inches(0.5), Inches(0.95),
        SLIDE_W - Inches(1.0), Inches(0.6),
        title_clean, font_size=26, bold=True, color=BRAND_DARK
    )

    # Accent underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.55),
        Inches(1.2), Inches(0.05)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_CYAN

    # Body
    y = Inches(1.75)
    content_left = Inches(0.5)
    content_w = SLIDE_W - Inches(1.0)
    max_y = SLIDE_H - Inches(0.8)

    for kind, content in slide_data.blocks:
        if y > max_y:
            break
        if kind == "para":
            text = clean_md(content)
            # Count lines roughly for height
            h = min(Inches(3.0), Inches(0.45 + 0.04 * len(text)))
            tb = slide.shapes.add_textbox(content_left, y, content_w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.text = ""
            # Support inline bold
            parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", content)
            for part in parts:
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**"):
                    r.text = part[2:-2]
                    r.font.bold = True
                elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                    r.text = part[1:-1]
                    r.font.italic = True
                else:
                    r.text = part
                r.font.size = Pt(16)
                r.font.color.rgb = BRAND_DARKER
                r.font.name = "Calibri"
            y += h + Inches(0.1)

        elif kind == "bullet":
            bullets = content
            # Estimate height
            h = Inches(0.45 * len(bullets) + 0.2)
            if y + h > max_y:
                h = max_y - y
            tb = slide.shapes.add_textbox(content_left, y, content_w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, bt in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = ""
                p.space_after = Pt(4)
                # Bullet marker
                rb = p.add_run()
                rb.text = "▸ "
                rb.font.size = Pt(16)
                rb.font.color.rgb = BRAND_BLUE
                rb.font.bold = True
                rb.font.name = "Calibri"
                # Bullet content (inline bold)
                parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", bt)
                for part in parts:
                    if not part:
                        continue
                    r = p.add_run()
                    if part.startswith("**") and part.endswith("**"):
                        r.text = part[2:-2]
                        r.font.bold = True
                        r.font.color.rgb = BRAND_DARK
                    elif part.startswith("*") and part.endswith("*") \
                            and len(part) > 2:
                        r.text = part[1:-1]
                        r.font.italic = True
                        r.font.color.rgb = BRAND_DARKER
                    else:
                        r.text = part
                        r.font.color.rgb = BRAND_DARKER
                    r.font.size = Pt(15)
                    r.font.name = "Calibri"
            y += h + Inches(0.1)

        elif kind == "table":
            rows = parse_table(content)
            if not rows:
                continue
            n_rows = len(rows)
            n_cols = len(rows[0])
            row_h = Inches(min(0.5, max(0.28, 5.0 / max(n_rows, 1))))
            table_h = row_h * n_rows
            if y + table_h > max_y:
                table_h = max_y - y - Inches(0.2)
                row_h = table_h / n_rows
            gtable = slide.shapes.add_table(
                n_rows, n_cols,
                content_left, y, content_w, table_h
            ).table

            for i, row in enumerate(rows):
                for j, cell_text in enumerate(row):
                    if j >= n_cols:
                        break
                    cell = gtable.cell(i, j)
                    cell.text = ""
                    cell.fill.solid()
                    if i == 0:
                        cell.fill.fore_color.rgb = BRAND_BLUE
                        text_color = BRAND_WHITE
                        is_bold = True
                    else:
                        if i % 2 == 1:
                            cell.fill.fore_color.rgb = BRAND_WHITE
                        else:
                            cell.fill.fore_color.rgb = BRAND_LIGHT
                        text_color = BRAND_DARKER
                        is_bold = False
                    tf = cell.text_frame
                    tf.margin_left = Inches(0.08)
                    tf.margin_right = Inches(0.08)
                    tf.margin_top = Inches(0.02)
                    tf.margin_bottom = Inches(0.02)
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = ""
                    p.alignment = PP_ALIGN.CENTER if j == 0 or \
                        re.match(r"^[\d\-\.,%⭐✅❌🟡🟠⚪⭐ ]+$", clean_md(cell_text)) \
                        else PP_ALIGN.LEFT
                    parts = re.split(r"(\*\*[^*]+\*\*)", cell_text)
                    for part in parts:
                        if not part:
                            continue
                        r = p.add_run()
                        if part.startswith("**") and part.endswith("**"):
                            r.text = part[2:-2]
                            r.font.bold = True
                        else:
                            r.text = part
                            r.font.bold = is_bold
                        font_size = 11 if n_rows > 6 else 12
                        if n_rows > 9:
                            font_size = 10
                        r.font.size = Pt(font_size)
                        r.font.color.rgb = text_color
                        r.font.name = "Calibri"

            y += table_h + Inches(0.15)

        elif kind == "quote":
            text = clean_md(content)
            h = Inches(max(0.6, 0.4 + 0.03 * len(text)))
            # Quote box with accent
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                content_left, y, content_w, h
            )
            box.line.fill.background()
            box.fill.solid()
            box.fill.fore_color.rgb = BRAND_LIGHT
            # Left accent strip inside
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                content_left + Inches(0.05), y + Inches(0.1),
                Inches(0.1), h - Inches(0.2)
            )
            strip.line.fill.background()
            strip.fill.solid()
            strip.fill.fore_color.rgb = BRAND_BLUE

            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.35)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = ""
            parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", content)
            for part in parts:
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**"):
                    r.text = part[2:-2]
                    r.font.bold = True
                elif part.startswith("*") and part.endswith("*") \
                        and len(part) > 2:
                    r.text = part[1:-1]
                    r.font.italic = True
                else:
                    r.text = part
                r.font.size = Pt(15)
                r.font.italic = True
                r.font.color.rgb = BRAND_DARK
                r.font.name = "Calibri"
            y += h + Inches(0.1)


# ============================================================
# 5. BUILD PPTX FROM MARKDOWN
# ============================================================
def build_pptx(md_path: Path, out_path: Path, sections: list[tuple[int, str]]):
    md_text = md_path.read_text(encoding="utf-8")
    meta, slides = parse_markdown(md_text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title = meta.get("title", "TNTech")
    subtitle = meta.get("subtitle", "")
    date_text = meta.get("date", "")

    # --- Cover ---
    build_cover_slide(prs, title, subtitle, date_text)

    # --- Separate first slide (is_cover from H1) and agenda ---
    # First slide is cover (H1) - skip it since we rebuilt from meta
    # Second slide is agenda (## Nội dung trình bày)
    if slides and getattr(slides[0], "_is_cover", False):
        slides = slides[1:]

    # Total count for pagination: cover + content slides + thank you
    # Agenda = slides[0], content = slides[1:-1], thank-you = slides[-1]
    total = len(prs.slides._sldIdLst) + len(slides)  # will recalc

    # Pre-compute inserted sections
    section_positions = {sec_after_slide: (num, title)
                         for num, title, sec_after_slide in
                         [(s[0], s[1], s[2]) for s in sections]
                         } if sections else {}

    # Walk through slides: agenda, then content slides, then last slide = thanks
    # Agenda
    if slides:
        first = slides[0]
        # Build agenda items from bullet block
        items = []
        for kind, content in first.blocks:
            if kind == "bullet":
                items = content
                break
        if items:
            # placeholder slide_num=2, total=?  - we'll recount later
            build_agenda_slide(prs, items, slide_num=2, total_slides=0)
        else:
            build_content_slide(prs, first, 2, 0)
        slides = slides[1:]

    # Detect thank-you last slide (title starts with "Cảm ơn")
    last_slide = None
    if slides and ("Cảm ơn" in slides[-1].title or "cảm ơn" in slides[-1].title.lower()):
        last_slide = slides[-1]
        slides = slides[:-1]

    # Content slides
    for i, sl in enumerate(slides):
        # Insert section divider if configured
        if sections:
            for sec_num, sec_title, after_slide_idx in sections:
                if i == after_slide_idx:
                    build_section_divider(prs, sec_num, sec_title, 0, 0)
        build_content_slide(prs, sl, 0, 0)

    # Thank you
    if last_slide:
        # Extract contact bullets
        bullets = []
        for kind, content in last_slide.blocks:
            if kind == "bullet":
                bullets.extend([clean_md(b) for b in content])
            elif kind == "para":
                t = clean_md(content).strip()
                if t and len(t) < 200:
                    bullets.append(t)
        build_thankyou_slide(prs, "Cảm ơn Quý khách!", bullets[:6])

    # Recount and update slide numbers
    total_slides = len(prs.slides)
    # Update footer page numbers
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Match "N / 0" pattern from our placeholder
                        m = re.match(r"^\d+\s*/\s*0$", run.text.strip())
                        if m:
                            run.text = f"{idx} / {total_slides}"

    prs.save(str(out_path))
    return total_slides


# ============================================================
# 6. MAIN
# ============================================================
if __name__ == "__main__":
    base = Path(__file__).parent

    # ENTERPRISE
    ent_md = base / "PROPOSAL-ERP-ENTERPRISE-SLIDES.md"
    ent_pptx = base / "PROPOSAL-ERP-ENTERPRISE-TNTech.pptx"
    n1 = build_pptx(ent_md, ent_pptx, sections=None)
    print(f"[OK] {ent_pptx.name} — {n1} slides")

    # HOSPITAL
    hosp_md = base / "PROPOSAL-ERP-HOSPITAL-SLIDES.md"
    hosp_pptx = base / "PROPOSAL-ERP-HOSPITAL-TNTech.pptx"
    n2 = build_pptx(hosp_md, hosp_pptx, sections=None)
    print(f"[OK] {hosp_pptx.name} — {n2} slides")

    print(f"\nLogo saved: {LOGO_PATH}")
    print("Done.")
